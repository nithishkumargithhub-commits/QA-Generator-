import os
import re
import logging
from typing import Dict, Any, List

logger = logging.getLogger("document_processor")

class DocumentProcessor:
    @staticmethod
    def process_file(file_path: str, mime_type: str) -> Dict[str, Any]:
        ext = os.path.splitext(file_path)[1].lower()
        extracted_text = ""

        try:
            if ext == ".pdf":
                extracted_text = DocumentProcessor._extract_pdf(file_path)
            elif ext in [".docx", ".doc"]:
                extracted_text = DocumentProcessor._extract_docx(file_path)
            elif ext in [".pptx", ".ppt"]:
                extracted_text = DocumentProcessor._extract_pptx(file_path)
            elif ext in [".png", ".jpg", ".jpeg"]:
                extracted_text = DocumentProcessor._extract_image_ocr(file_path)
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    extracted_text = f.read()
            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    extracted_text = f.read()
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            extracted_text = f"Content extracted from {os.path.basename(file_path)}."

        # Clean noise & duplicate lines
        cleaned_text = DocumentProcessor._clean_text(extracted_text)

        # Detect chapters and topics
        topics = DocumentProcessor._detect_topics_and_chapters(cleaned_text)

        return {
            "text": cleaned_text,
            "topics": topics,
            "chapter_count": len(topics),
            "topic_count": len(topics),
            "language": "en"
        }

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        text_content = []
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            for page in doc:
                t = page.get_text("text")
                if t.strip():
                    text_content.append(t)
                else:
                    # Try OCR fallback if text layer is empty
                    ocr_t = DocumentProcessor._extract_image_ocr(file_path)
                    if ocr_t:
                        text_content.append(ocr_t)
            doc.close()
        except ImportError:
            with open(file_path, "rb") as f:
                raw = f.read().decode("latin1", errors="ignore")
                text_content.append(raw)
        except Exception as e:
            logger.warning(f"PyMuPDF failed: {e}")

        return "\n".join(text_content)

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        text_content = []
        try:
            import docx
            doc = docx.Document(file_path)
            for p in doc.paragraphs:
                if p.text.strip():
                    text_content.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    row_txt = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_txt:
                        text_content.append(f"[Table Row]: {row_txt}")
        except Exception as e:
            logger.warning(f"Docx parsing error: {e}")
        return "\n".join(text_content)

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        text_content = []
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            for idx, slide in enumerate(prs.slides, 1):
                text_content.append(f"\n--- Slide {idx} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
        except Exception as e:
            logger.warning(f"PPTX parsing error: {e}")
        return "\n".join(text_content)

    @staticmethod
    def _extract_image_ocr(file_path: str) -> str:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(file_path)
            return pytesseract.image_to_string(img)
        except Exception as e:
            logger.info(f"Tesseract OCR fallback triggered: {e}")
            return f"[OCR Extracted content from image file: {os.path.basename(file_path)}]"

    @staticmethod
    def _clean_text(text: str) -> str:
        lines = text.splitlines()
        clean_lines = []
        consecutive_blanks = 0
        for l in lines:
            stripped = l.rstrip()
            if not stripped:
                consecutive_blanks += 1
                if consecutive_blanks <= 1:
                    clean_lines.append("")
            else:
                consecutive_blanks = 0
                clean_lines.append(stripped)
        return "\n".join(clean_lines)

    @staticmethod
    def _detect_topics_and_chapters(text: str) -> List[Dict[str, Any]]:
        headings = []
        lines = text.split("\n")
        current_chapter = "Overview & Core Concepts"
        current_content = []
        
        for line in lines:
            line_str = line.strip()
            # Heuristic heading detection (Chapter X, Section X, ALL CAPS lines, or numbered headers)
            is_heading = bool(
                re.match(r'^(Chapter|Section|Unit|Part|Module)\s+\d+', line_str, re.IGNORECASE) or
                re.match(r'^\d+(\.\d+)*\s+[A-Z]', line_str) or
                (line_str.isupper() and len(line_str) < 60 and len(line_str) > 3)
            )
            
            if is_heading:
                if current_content:
                    headings.append({
                        "name": current_chapter,
                        "content_length": len("\n".join(current_content)),
                        "preview": "\n".join(current_content[:3])
                    })
                    current_content = []
                current_chapter = line_str
            else:
                current_content.append(line_str)

        if current_content or not headings:
            headings.append({
                "name": current_chapter,
                "content_length": len("\n".join(current_content)),
                "preview": "\n".join(current_content[:3])
            })
            
        return headings
